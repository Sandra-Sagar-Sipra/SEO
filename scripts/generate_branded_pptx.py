"""
Generate a branded PowerPoint (.pptx) from embedded slide data.
Requires: pip install -r scripts/requirements-pptx.txt

Logos (repo root, exact filenames):
  - SipraHub: assets/logos/siprahub-logo.png
  - Aiquire:  assets/logos/Aiquire_logo.svg (SVG is converted via cairosvg; pip install cairosvg)
"""

from __future__ import annotations

import argparse
import io
import json
import re
import sys
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

# Repo root = parent of scripts/
ROOT = Path(__file__).resolve().parent.parent

DEFAULT_DECK_JSON = ROOT / "output" / "decks" / "ai-in-business-workflows-siprahub.json"

# Fallback if default JSON is missing (minimal generic deck).
_FALLBACK_JSON = r"""{
  "brand": "siprahub",
  "slides": [{"title": "Branded deck", "points": ["Add output/decks/*.json or pass a path.", "Place logos under assets/logos/."]}]
}"""


def _strip_frontmatter(text: str) -> str:
    if text.startswith("---"):
        parts = text.split("---", 2)
        if len(parts) >= 3:
            return parts[2]
    return text


def _strip_md_links(s: str) -> str:
    return re.sub(r"\[([^\]]+)\]\([^)]*\)", r"\1", s)


def _prose_to_bullets(chunk: str, max_points: int = 5) -> list[str]:
    chunk = _strip_md_links(chunk)
    chunk = re.sub(r"\*\*([^*]+)\*\*", r"\1", chunk)
    out: list[str] = []
    for block in re.split(r"\n\s*\n+", chunk.strip()):
        if not block.strip() or block.strip().startswith(">"):
            continue
        sentences = re.split(r"(?<=[.!?])\s+", block.replace("\n", " "))
        for sent in sentences:
            sent = sent.strip()
            if len(sent) < 25:
                continue
            if sent in out:
                continue
            sent = sent[:300] + ("…" if len(sent) > 300 else "")
            out.append(sent)
            if len(out) >= max_points:
                return out
    return out[:max_points]


def _body_to_points(body: str) -> list[str]:
    body = _strip_md_links(body)
    points: list[str] = []
    for line in body.splitlines():
        line = line.strip()
        if not line or line.startswith(">"):
            continue
        m = re.match(r"^[\-*]\s+(.+)", line)
        if m:
            points.append(m.group(1).strip())
            continue
        m = re.match(r"^\d+\.\s+\*?\*?(.+?)\*?\*?:?\s*$", line)
        if m and not line.startswith("```"):
            points.append(m.group(1).strip().strip("*"))
    if len(points) >= 2:
        return points[:5]
    return _prose_to_bullets(body, max_points=5)


def markdown_article_to_payload(md_path: Path, brand: str) -> dict:
    """Build deck payload from a blog article: H1 + H2 sections become slides (skips Meta/SEO/Engagement)."""
    text = _strip_frontmatter(md_path.read_text(encoding="utf-8"))
    slides: list[dict] = []

    sections = re.split(r"^## ", text, flags=re.MULTILINE)
    head = sections[0].strip()
    for line in head.splitlines():
        if line.startswith("# "):
            h1 = line[2:].strip()
            rest = head.split(line, 1)[1]
            rest = rest.split("---", 1)[0].strip()
            intro_pts = _prose_to_bullets(rest, max_points=5)
            if h1 and intro_pts:
                slides.append({"title": h1, "points": intro_pts})
            break

    skip_substrings = ("meta elements", "seo checklist", "engagement checklist")
    for block in sections[1:]:
        if not block.strip():
            continue
        lines = block.splitlines()
        heading = lines[0].strip()
        low = heading.lower()
        if any(s in low for s in skip_substrings):
            break
        body = "\n".join(lines[1:])
        if "```" in body:
            body = body.split("```", 1)[0]
        pts = _body_to_points(body)
        if heading and pts:
            slides.append({"title": heading, "points": pts})

    return {"brand": brand, "slides": slides}


def load_payload(path: Path | None) -> dict:
    if path and path.is_file():
        return json.loads(path.read_text(encoding="utf-8"))
    if DEFAULT_DECK_JSON.is_file():
        return json.loads(DEFAULT_DECK_JSON.read_text(encoding="utf-8"))
    return json.loads(_FALLBACK_JSON)


def logo_path_for_brand(brand: str) -> Path:
    b = (brand or "").strip().lower()
    if b == "siprahub":
        return ROOT / "assets" / "logos" / "siprahub-logo.png"
    if b == "aiquire":
        return ROOT / "assets" / "logos" / "Aiquire_logo.svg"
    raise ValueError(f"Unknown brand: {brand!r}. Use 'siprahub' or 'aiquire'.")


def _add_logo_picture(slide, logo_path: Path) -> None:
    """python-pptx needs raster images; SVG is converted to PNG in memory."""
    left, top, width = Inches(0.3), Inches(0.2), Inches(1)
    if logo_path.suffix.lower() == ".svg":
        try:
            import cairosvg
        except ImportError as err:
            raise ImportError(
                "Aiquire logo is SVG. Install: pip install cairosvg\n"
                "Or add a raster copy as assets/logos/aiquire.png and edit logo_path_for_brand."
            ) from err
        png_bytes = cairosvg.svg2png(url=str(logo_path))
        stream = io.BytesIO(png_bytes)
        slide.shapes.add_picture(stream, left, top, width=width)
    else:
        slide.shapes.add_picture(str(logo_path), left, top, width=width)


def build_presentation(payload: dict) -> Presentation:
    brand = payload["brand"]
    slides_spec = payload["slides"]
    logo = logo_path_for_brand(brand)
    if not logo.is_file():
        raise FileNotFoundError(
            f"Logo not found: {logo}. Add the file at this path (see script docstring)."
        )

    prs = Presentation()
    # Title and content (layout index 1 on default template)
    layout = prs.slide_layouts[1]

    for item in slides_spec:
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = item["title"]
        tf = slide.placeholders[1].text_frame
        pts = item["points"]
        if pts:
            tf.paragraphs[0].text = pts[0]
            for point in pts[1:]:
                p = tf.add_paragraph()
                p.text = point
                p.level = 0
        _add_logo_picture(slide, logo)

    return prs


def main() -> None:
    parser = argparse.ArgumentParser(description="Generate branded .pptx from JSON or from a markdown article.")
    parser.add_argument(
        "json",
        nargs="?",
        default=None,
        help="Path to deck JSON ({brand, slides:[{title, points}]}). If omitted, uses default deck JSON.",
    )
    parser.add_argument(
        "--from-md",
        dest="from_md",
        type=Path,
        default=None,
        help="Build slides from a markdown blog file (H1 intro + ## sections; skips Meta/SEO checklists).",
    )
    parser.add_argument(
        "--brand",
        default=None,
        help="Override brand (siprahub|aiquire) when using --from-md.",
    )
    args = parser.parse_args()

    if args.from_md:
        brand = (args.brand or "siprahub").strip().lower()
        payload = markdown_article_to_payload(args.from_md.resolve(), brand)
    else:
        json_path = Path(args.json).resolve() if args.json else None
        payload = load_payload(json_path)
        if args.brand:
            payload = {**payload, "brand": args.brand.strip().lower()}

    brand = payload["brand"]
    prs = build_presentation(payload)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    out_name = f"{brand}_presentation_{timestamp}.pptx"
    out_path = ROOT / "output" / out_name
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(out_path)


if __name__ == "__main__":
    main()
